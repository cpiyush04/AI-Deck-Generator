import google.generativeai as genai
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
import json
import io
import urllib.parse
from PIL import Image
from dotenv import load_dotenv
import os

# --- Load Environment Variables ---
load_dotenv()

# --- Configuration ---
# Keys are now loaded securely from the .env file
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
GOOGLE_SEARCH_API_KEY = os.getenv("GOOGLE_SEARCH_API_KEY")
CUSTOM_SEARCH_ENGINE_ID = os.getenv("CUSTOM_SEARCH_ENGINE_ID")

if "YOUR_GEMINI_API_KEY" in GOOGLE_API_KEY:
    print("ERROR: Please set your GOOGLE_API_KEY in the script.")
    exit()

genai.configure(api_key=GOOGLE_API_KEY)

# --- Agent Definitions ---

class ResearchAgent:
    """An agent that scours the web for information on a topic using multiple search engines."""
    
    def _search_duckduckgo(self, topic):
        """Performs a search on DuckDuckGo and returns context."""
        try:
            print("  -> Searching DuckDuckGo...")
            search_url = f"https://html.duckduckgo.com/html/?q={urllib.parse.quote(topic)}"
            headers = {'User-Agent': 'Mozilla/5.0'}
            response = requests.get(search_url, headers=headers, timeout=10)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, 'html.parser')
            snippets = [p.get_text() for p in soup.find_all('a', class_='result__snippet')]
            return " ".join(snippets[:3]) if snippets else ""
        except Exception as e:
            print(f"      -> DuckDuckGo search failed: {e}")
            return ""

    def _search_google(self, topic):
        """Performs a search using Google's Custom Search API."""
        if "YOUR_GOOGLE_SEARCH_API_KEY" in GOOGLE_SEARCH_API_KEY or "YOUR_CSE_ID" in CUSTOM_SEARCH_ENGINE_ID:
            print("  -> Google Search keys not provided. Skipping Google search.")
            return ""
        try:
            print("  -> Searching Google...")
            api_url = f"https://www.googleapis.com/customsearch/v1?key={GOOGLE_SEARCH_API_KEY}&cx={CUSTOM_SEARCH_ENGINE_ID}&q={urllib.parse.quote(topic)}"
            response = requests.get(api_url, timeout=10)
            response.raise_for_status()
            search_results = response.json()
            snippets = [item.get('snippet', '') for item in search_results.get('items', [])]
            return " ".join(snippets[:3]) if snippets else ""
        except Exception as e:
            print(f"      -> Google search failed: {e}")
            return ""

    def search_web_for_topic(self, topic):
        print(f"🕵️  ResearchAgent: Searching the web for '{topic}'...")
        ddg_context = self._search_duckduckgo(topic)
        google_context = self._search_google(topic)
        
        combined_context = (ddg_context + " " + google_context).strip()
        
        if combined_context:
            print("  -> Found relevant context from web search.")
            return combined_context
        else:
            print("  -> Web search failed to find context.")
            return "No specific web context found."


class OrchestratorAgent:
    """An agent that provides a fixed 7-slide structure for the presentation."""
    def create_presentation_plan(self, topic, web_context):
        print("🧠 OrchestratorAgent: Providing fixed 7-slide presentation plan...")
        
        plan = {
            "slides": [
                {"type": "title_slide", "purpose": "A compelling title for the presentation."},
                {"type": "overview_slide", "purpose": "An overview describing key talking points."},
                {"type": "key_point_slide", "purpose": f"The first key point or trend about {topic}."},
                {"type": "key_point_slide", "purpose": f"The second key point or argument about {topic}."},
                {"type": "key_point_slide", "purpose": f"The third key point or trend about {topic}."},
                {"type": "key_point_slide", "purpose": f"The fourth key point or argument about {topic}."},
                {"type": "conclusion_slide", "purpose": "Give Concluding Points."}
            ]
        }
        print("  -> Fixed plan provided successfully.")
        return plan

class ContentStrategistAgent:
    """An agent that executes a plan to write the presentation content."""
    def generate_presentation_content(self, topic, web_context, plan):
        print("✍️  ContentStrategistAgent: Executing plan and writing content...")
        model = genai.GenerativeModel('gemini-2.5-flash')
        
        plan_str = json.dumps(plan)
        prompt = f"""
        Act as a Content Strategist. Your task is to generate the text content for a presentation about "{topic}".
        Use the provided web context: "{web_context}".
        You MUST follow this presentation plan exactly: {plan_str}.

        For each slide in the plan, generate a "title" and a list of "points".
        *** IMPORTANT: Generate only 4 or 5 bullet points per slide. Each point must be a short but explanatory sentence, not just a heading. ***
        
        Return a single, raw JSON object that contains a "slides" key, which is a list of slide objects with the content filled in.
        
        Example output structure:
        {{
            "slides": [
                {{"type": "title_slide", "title": "The Evolution of Artificial Intelligence", "points": []}},
                {{"type": "overview_slide", "title": "Overview", "points": ["Exploring the foundational concepts of AI.", "Covering the major milestones and breakthroughs."]}},
                {{"type": "key_point_slide", "title": "Early Concepts and the Turing Test", "points": ["Alan Turing's 1950 paper proposed a test for machine intelligence, now known as the Turing Test.", "Early research focused on problem-solving and symbolic methods."]}},
                {{"type": "conclusion_slide", "title": "Conclusion", "points": ["AI has evolved from simple concepts to complex neural networks.", "The future of AI holds immense potential for innovation across all industries."]}}
            ]
        }}
        """
        try:
            response = model.generate_content(prompt)
            # Clean up the response to ensure it's valid JSON
            content_json = response.text.strip().replace("```json", "").replace("```", "")
            print("  -> AI content structure generated successfully.")
            return json.loads(content_json)
        except Exception as e:
            print(f"  -> AI content generation failed: {e}")
            return None

class VisualAssetAgent:
    """An agent that determines the best visual assets for the content."""
    def add_image_queries_to_content(self, content):
        print("🎨 VisualAssetAgent: Determining image queries...")
        model = genai.GenerativeModel('gemini-2.5-flash')
        
        content_str = json.dumps(content)
        prompt = f"""
        Act as a Visual Asset Director. For each slide with a "type" of "key_point_slide" in the following JSON object, add a simple, one or two-word "image_query".
        This query should be perfect for finding a high-quality stock photo on a site like Pixabay or Google Images.
        Return the entire JSON object, now including the "image_query" fields where appropriate.
        
        JSON Input:
        {content_str}
        
        Return nothing but the updated, raw JSON object.
        """
        try:
            response = model.generate_content(prompt)
            enriched_content_json = response.text.strip().replace("```json", "").replace("```", "")
            print("  -> Image queries added successfully.")
            return json.loads(enriched_content_json)
        except Exception as e:
            print(f"  -> Failed to add image queries: {e}")
            # Return original content if the AI fails
            return content

class PowerPointAssemblyAgent:
    """An agent that assembles the final PowerPoint presentation."""
    def _get_image_for_slide(self, query):
        """Searches for an image and converts it to a supported format (PNG)."""
        if "YOUR_GOOGLE_SEARCH_API_KEY" in GOOGLE_SEARCH_API_KEY or "YOUR_CSE_ID" in CUSTOM_SEARCH_ENGINE_ID:
            print("  -> Google Search keys not provided. Skipping image insertion.")
            return None
            
        print(f"🖼️  AssemblyAgent: Searching Google Images for '{query}'...")
        try:
            api_url = f"https://www.googleapis.com/customsearch/v1?key={GOOGLE_SEARCH_API_KEY}&cx={CUSTOM_SEARCH_ENGINE_ID}&q={urllib.parse.quote(query)}&searchType=image&num=1"
            response = requests.get(api_url, timeout=10)
            response.raise_for_status()
            search_results = response.json()
            
            if 'items' not in search_results or not search_results['items']:
                print("  -> No image found on Google.")
                return None

            image_url = search_results['items'][0]['link']
            
            image_response = requests.get(image_url, timeout=10)
            image_response.raise_for_status()
            
            print("  -> Google image downloaded successfully. Converting to PNG...")

            # Open the downloaded image data with Pillow
            image = Image.open(io.BytesIO(image_response.content))
            
            # Create a new in-memory byte stream to hold the converted image
            output_stream = io.BytesIO()
            
            # Save the image to the stream in PNG format
            image.save(output_stream, format='PNG')
            output_stream.seek(0) # Rewind the stream to the beginning
            
            print("  -> Image converted successfully.")
            return output_stream

        except Exception as e:
            print(f"  -> Google Image search/download/conversion failed: {e}")
            return None


    def create_powerpoint_deck(self, content, topic):
        print("🛠️  PowerPointAssemblyAgent: Building the presentation...")
        prs = Presentation()
        
        if 'slides' not in content or not isinstance(content['slides'], list):
            print("  -> ERROR: Content is not in the expected format with a 'slides' list.")
            return None

        for slide_data in content['slides']:
            slide_type = slide_data.get('type')
            title = slide_data.get('title', 'Untitled Slide')
            points = slide_data.get('points', [])

            if slide_type == 'title_slide':
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                slide.shapes.title.text = title
                slide.placeholders[1].text = f"A Presentation on {topic}"
            
            elif slide_type in ['overview_slide', 'conclusion_slide']:
                slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
                slide.shapes.title.text = title
                tf = slide.shapes.placeholders[1].text_frame
                tf.clear()
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                
                for point in points:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0
                    p.space_after = Pt(8) 
                    p.font.size = Pt(15)

            elif slide_type == 'key_point_slide':
                image_query = slide_data.get("image_query", topic)
                image_stream = self._get_image_for_slide(image_query)

                # If an image is found, use the two-content layout.
                if image_stream:
                    slide = prs.slides.add_slide(prs.slide_layouts[3]) # Two Content
                    slide.shapes.title.text = title
                    
                    # Add text content to the left placeholder
                    left_placeholder = slide.placeholders[1]
                    tf = left_placeholder.text_frame
                    tf.clear()
                    tf.word_wrap = True
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    for point in points:
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
                        p.space_after = Pt(8)
                        p.font.size = Pt(18)
                    
                    # Add image to the right placeholder
                    right_placeholder = slide.placeholders[2]
                    slide.shapes.add_picture(
                        image_stream,
                        right_placeholder.left,
                        right_placeholder.top,
                        width=right_placeholder.width,
                        height=right_placeholder.height
                    )
                # If no image is found, use the standard Title and Content layout.
                else:
                    slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title and Content
                    slide.shapes.title.text = title
                    tf = slide.shapes.placeholders[1].text_frame
                    tf.clear()
                    tf.word_wrap = True
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    
                    for point in points:
                        p = tf.add_paragraph()
                        p.text = point
                        p.level = 0
                        p.space_after = Pt(8) 
                        p.font.size = Pt(18)
            
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        print("  -> Presentation assembled successfully.")
        return file_stream


# --- Main Execution Block ---
if __name__ == '__main__':
    # ---GET THE PRESENTATION TOPIC FROM THE USER ---
    presentation_topic = input("Please enter the topic for your presentation: ")

    print(f"\n--- Starting Presentation Generation for: '{presentation_topic}' ---")

    # ---ORCHESTRATE THE AGENTIC WORKFLOW ---
    researcher = ResearchAgent()
    orchestrator = OrchestratorAgent()
    content_strategist = ContentStrategistAgent()
    visual_director = VisualAssetAgent()
    assembler = PowerPointAssemblyAgent()

    # Research the topic
    web_context = researcher.search_web_for_topic(presentation_topic)
    
    # Create a plan
    presentation_plan = orchestrator.create_presentation_plan(presentation_topic, web_context)
    if not presentation_plan:
        print("ERROR: Failed to create presentation plan. Exiting.")
        exit()

    # Generate the initial text content
    initial_content = content_strategist.generate_presentation_content(presentation_topic, web_context, presentation_plan)
    if not initial_content:
        print("ERROR: Failed to generate initial content. Exiting.")
        exit()
        
    # Add image search queries to the content
    enriched_content = visual_director.add_image_queries_to_content(initial_content)
    
    # Build the PowerPoint file
    ppt_stream = assembler.create_powerpoint_deck(enriched_content, presentation_topic)
    
    # ---SAVE THE FINAL PRESENTATION ---
    if ppt_stream:
        output_filename = f"presentation-{presentation_topic}.pptx"
        try:
            with open(output_filename, 'wb') as f:
                f.write(ppt_stream.getbuffer())
            print(f"\nWorkflow Complete. Presentation saved as '{output_filename}'")
        except Exception as e:
            print(f"\nERROR: Failed to save the presentation file: {e}")
    else:
        print("\nERROR: Failed to assemble the PowerPoint file. No file was saved.")